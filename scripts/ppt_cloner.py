#!/usr/bin/env python3
"""
ppt_cloner.py - 使用 python-pptx 实现可靠的 PPT 克隆

核心思路：
1. 打开模板 PPT
2. 删除不需要的幻灯片
3. 复制需要的幻灯片（通过复制+删除原始的方式）
4. 替换文本内容

这种方式让 python-pptx 自动处理所有内部关系，避免"内容有问题"弹窗
"""

import json
import copy
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from typing import Dict, List, Optional
from collections import defaultdict
import re


def emu_to_pt(emu: int) -> float:
    """EMU 转换为 points"""
    return round(emu / 914400 * 72, 2) if emu else 0


class SlideInfo:
    """幻灯片信息"""
    def __init__(self, index: int):
        self.index = index
        self.type = 'content'
        self.layout_name = ''
        self.text_elements: List[Dict] = []
        self.shape_count = 0
        self.preview_text = ''
    
    def to_dict(self) -> Dict:
        return {
            'index': self.index,
            'type': self.type,
            'layout_name': self.layout_name,
            'text_elements': self.text_elements,
            'shape_count': self.shape_count,
            'preview_text': self.preview_text
        }


class PPTCloner:
    """PPT 克隆器 - 使用 python-pptx 实现可靠克隆"""
    
    def __init__(self, template_path: str):
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        
        # 分析模板
        self.analysis = self._analyze_template()
    
    def _analyze_template(self) -> Dict:
        """分析模板结构"""
        prs = Presentation(str(self.template_path))
        
        result = {
            'source': self.template_path.name,
            'slide_count': len(prs.slides),
            'slide_width_pt': emu_to_pt(prs.slide_width),
            'slide_height_pt': emu_to_pt(prs.slide_height),
            'slides': [],
            'slide_types': defaultdict(list),
        }
        
        for idx, slide in enumerate(prs.slides):
            info = self._analyze_slide(slide, idx, len(prs.slides))
            result['slides'].append(info.to_dict())
            result['slide_types'][info.type].append(idx)
        
        result['slide_types'] = dict(result['slide_types'])
        return result
    
    def _analyze_slide(self, slide, idx: int, total: int) -> SlideInfo:
        """分析单张幻灯片"""
        info = SlideInfo(idx)
        info.layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
        info.shape_count = len(slide.shapes)
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    info.text_elements.append({
                        'shape_name': shape.name,
                        'text': text,
                    })
                    if not info.preview_text:
                        info.preview_text = text[:50]
        
        info.type = self._infer_slide_type(info, idx, total)
        return info
    
    def _infer_slide_type(self, info: SlideInfo, idx: int, total: int) -> str:
        """推断幻灯片类型"""
        layout_lower = info.layout_name.lower()
        all_text = ' '.join([t['text'].lower() for t in info.text_elements])
        
        if idx == 0:
            return 'cover'
        
        if any(k in layout_lower for k in ['title', 'cover', '封面']):
            return 'cover'
        if any(k in layout_lower for k in ['section', 'divider', '章节']):
            return 'divider'
        
        if any(k in all_text for k in ['目录', 'contents', 'agenda']):
            return 'toc'
        if any(k in all_text for k in ['谢谢', 'thank', '感谢', '聆听']):
            return 'ending'
        
        if len(info.text_elements) <= 4:
            for t in info.text_elements:
                text = t['text'].strip()
                if re.match(r'^[0-9０-９]+$', text) or text in ['01', '02', '03', '04', '05']:
                    return 'divider'
        
        return 'content'
    
    def get_slides_by_type(self, slide_type: str) -> List[int]:
        """获取指定类型的幻灯片索引"""
        return self.analysis['slide_types'].get(slide_type, [])
    
    def create_from_plan(self, content_plan: List[Dict], output_path: str) -> str:
        """
        根据内容计划创建新 PPT
        
        content_plan 格式:
        [
            {
                "template_slide": 0,  # 使用模板的第几张 (0-indexed)
                "replacements": {
                    "原文本": "新文本",
                    ...
                }
            },
            ...
        ]
        """
        # 打开模板
        prs = Presentation(str(self.template_path))
        
        # 收集需要保留的幻灯片索引和对应的替换规则
        slides_to_keep = []
        for item in content_plan:
            template_idx = item.get('template_slide')
            if template_idx is None:
                slide_type = item.get('type', 'content')
                type_slides = self.get_slides_by_type(slide_type)
                template_idx = type_slides[0] if type_slides else 0
            
            slides_to_keep.append({
                'source_idx': template_idx,
                'replacements': item.get('replacements', {})
            })
        
        # 构建要删除的幻灯片索引列表
        total_slides = len(prs.slides)
        indices_to_delete = set(range(total_slides))
        
        # 对于要保留的幻灯片，从删除列表中移除
        for item in slides_to_keep:
            indices_to_delete.discard(item['source_idx'])
        
        # 按倒序删除不需要的幻灯片（避免索引变化问题）
        for idx in sorted(indices_to_delete, reverse=True):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
        
        # 建立新旧索引映射
        old_to_new = {}
        new_idx = 0
        for old_idx in range(total_slides):
            if old_idx not in indices_to_delete:
                old_to_new[old_idx] = new_idx
                new_idx += 1
        
        # 应用文本替换
        for item in slides_to_keep:
            source_idx = item['source_idx']
            if source_idx in old_to_new:
                new_idx = old_to_new[source_idx]
                slide = prs.slides[new_idx]
                self._apply_replacements(slide, item['replacements'])
        
        # 重新排序幻灯片（按 content_plan 的顺序）
        # python-pptx 不直接支持重排，所以我们需要换一种方式
        
        # 保存
        prs.save(output_path)
        print(f"✓ 已生成: {output_path}")
        return output_path
    
    def create_simple(self, slide_indices: List[int], replacements_list: List[Dict], output_path: str) -> str:
        """
        简化版创建：指定要保留的幻灯片索引和对应的替换规则
        
        slide_indices: [0, 1, 2, 23]  # 保留第1、2、3、24张
        replacements_list: [{...}, {...}, {...}, {...}]  # 每张的替换规则
        """
        prs = Presentation(str(self.template_path))
        total = len(prs.slides)
        
        # 要删除的索引
        to_delete = set(range(total)) - set(slide_indices)
        
        # 按倒序删除
        for idx in sorted(to_delete, reverse=True):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
        
        # 建立索引映射
        old_to_new = {}
        new_idx = 0
        for old_idx in sorted(slide_indices):
            old_to_new[old_idx] = new_idx
            new_idx += 1
        
        # 应用替换
        for i, old_idx in enumerate(slide_indices):
            if old_idx in old_to_new and i < len(replacements_list):
                new_idx = old_to_new[old_idx]
                slide = prs.slides[new_idx]
                self._apply_replacements(slide, replacements_list[i])
        
        prs.save(output_path)
        print(f"✓ 已生成: {output_path}")
        return output_path
    
    def _apply_replacements(self, slide, replacements: Dict[str, str]):
        """应用文本替换"""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # 1. 检查名称匹配 (格式: "shape:ShapeName": "New Content")
            shape_name_key = f"shape:{shape.name}"
            if shape_name_key in replacements:
                shape.text = replacements[shape_name_key]
                continue

            # 2. 检查全文匹配 (如果形状的完整文本等于某个键，则全部替换)
            full_text = shape.text_frame.text.strip()
            if full_text in replacements:
                shape.text = replacements[full_text]
                continue

            # 3. 检查段落/Run 级别的部分匹配 (现有的逻辑，但进行了增强)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        original = run.text
                        for old_text, new_text in replacements.items():
                            if old_text.startswith("shape:"): continue
                            if old_text in original:
                                run.text = original.replace(old_text, new_text)
                                original = run.text
    
    def print_analysis(self):
        """打印分析结果"""
        print("\n" + "=" * 70)
        print(f"📊 模板分析: {self.analysis['source']}")
        print("=" * 70)
        
        print(f"\n📐 尺寸: {self.analysis['slide_width_pt']}pt x {self.analysis['slide_height_pt']}pt")
        print(f"📄 幻灯片数量: {self.analysis['slide_count']}")
        
        print("\n📑 幻灯片类型分布:")
        for slide_type, indices in self.analysis['slide_types'].items():
            print(f"   {slide_type}: 幻灯片 {[i+1 for i in indices]}")
        
        print("\n📝 各幻灯片详情:")
        for slide in self.analysis['slides']:
            idx = slide['index']
            stype = slide['type']
            preview = slide['preview_text'][:40] if slide['preview_text'] else '(无文本)'
            
            print(f"\n   【{idx+1}】类型: {stype}")
            print(f"       预览: {preview}...")
            
            if slide['text_elements']:
                print("       可替换文本:")
                for te in slide['text_elements'][:5]:
                    text = te['text'][:50]
                    print(f"         - \"{text}\"")
                if len(slide['text_elements']) > 5:
                    print(f"         ... 还有 {len(slide['text_elements']) - 5} 个")
        
        print("\n" + "=" * 70)
    
    def export_analysis(self, output_path: str):
        """导出分析结果"""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.analysis, f, ensure_ascii=False, indent=2)
        print(f"✓ 分析结果已保存到: {output_path}")


def main():
    if len(sys.argv) < 2:
        print("PPT 克隆器 - 基于模板创建新 PPT（无弹窗版）")
        print()
        print("用法:")
        print("  分析模板:")
        print("    python ppt_cloner.py analyze <template.pptx> [output.json]")
        print()
        print("  创建新 PPT:")
        print("    python ppt_cloner.py create <template.pptx> <plan.json> <output.pptx>")
        print()
        print("内容计划 JSON 格式:")
        print('''
[
    {
        "template_slide": 0,
        "replacements": {
            "原文本": "新文本"
        }
    }
]
        ''')
        sys.exit(1)
    
    command = sys.argv[1]
    
    if command == 'analyze':
        if len(sys.argv) < 3:
            print("错误: 请提供模板文件路径")
            sys.exit(1)
        
        template_path = sys.argv[2]
        output_json = sys.argv[3] if len(sys.argv) > 3 else None
        
        cloner = PPTCloner(template_path)
        cloner.print_analysis()
        
        if output_json:
            cloner.export_analysis(output_json)
    
    elif command == 'create':
        if len(sys.argv) < 5:
            print("错误: 请提供模板路径、内容计划 JSON 和输出路径")
            sys.exit(1)
        
        template_path = sys.argv[2]
        plan_path = sys.argv[3]
        output_path = sys.argv[4]
        
        with open(plan_path, 'r', encoding='utf-8') as f:
            content_plan = json.load(f)
        
        cloner = PPTCloner(template_path)
        
        # 提取 slide_indices 和 replacements_list
        slide_indices = []
        replacements_list = []
        
        for item in content_plan:
            template_idx = item.get('template_slide')
            if template_idx is None:
                slide_type = item.get('type', 'content')
                type_slides = cloner.get_slides_by_type(slide_type)
                template_idx = type_slides[0] if type_slides else 0
            
            slide_indices.append(template_idx)
            replacements_list.append(item.get('replacements', {}))
        
        cloner.create_simple(slide_indices, replacements_list, output_path)
    
    else:
        print(f"未知命令: {command}")
        sys.exit(1)


if __name__ == '__main__':
    main()
